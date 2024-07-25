from pptx import Presentation
import pandas as pd
import docx
import base64
import os
from func.generate_audio_stream import generate_audio_stream

def read_file(current_model):
    if not os.path.exists('./file'):
        os.makedirs('./file')
    file_folder = './file/'
    print(current_model + "：请输入文件名，文件切记放在file文件夹中", end="\n")
    filename = file_folder + input()
    if filename.endswith('.docx') or filename.endswith('.doc'):
        doc = docx.Document(filename)
        return ' '.join([paragraph.text for paragraph in doc.paragraphs]), False
    elif filename.endswith('.txt'):
        with open(filename, "r", encoding="utf-8") as f:
            return f.read(), False
    elif filename.endswith('.xlsx') or filename.endswith('.xls'):
        df = pd.read_excel(filename)
        return df.to_string(index=False), False
    elif filename.endswith('.pptx'):
        prs = Presentation(filename)
        text = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text.append(shape.text)
        return ' '.join(text), False
    elif filename.endswith('.png') or filename.endswith('jpg'):
        with open(filename, "rb") as image_file:
            return base64.b64encode(image_file.read()).decode("utf-8"), True
    else:
        print("目前不支持该文件格式")
        return "", False
    
def file_parse(current_model_name):
    new_screenshot = False
    file_output, new_screenshot = read_file(current_model_name)
    if file_output == "":
        return "", ""
    base64_image = ""
    if new_screenshot:
        base64_image = file_output
        file_output = ""
    recognized_text = file_output
    print("请说出你要对这份文档或者图片进行什么操作")
    generate_audio_stream("请说出你要对这份文档或者图片进行什么操作")
    return recognized_text, base64_image
